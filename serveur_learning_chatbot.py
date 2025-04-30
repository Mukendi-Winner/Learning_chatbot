import os
import json
from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
from together import Together
import PyPDF2
from pptx import Presentation
from werkzeug.utils import secure_filename
import uuid
from threading import Lock
import logging
import re

app = Flask(__name__)
CORS(app)

# Configuration du logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration
UPLOAD_FOLDER = 'Uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}
SESSION_STORAGE = {}
LOCK = Lock()

# Configurer Together AI
api_key = ""
client = Together(api_key=api_key)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_pdf_text(filepath):
    try:
        with open(filepath, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            pages = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ''
                logger.info(f"Page {i+1} extraite, longueur: {len(page_text)} caractères")
                sections = detect_sections(page_text)
                pages.append({
                    'page_number': i + 1,
                    'text': page_text,
                    'sections': sections
                })
            if not any(page['text'].strip() for page in pages):
                logger.warning("Aucun texte extrait du PDF")
                return "Aucun texte lisible n'a été extrait du PDF."
            logger.info(f"Nombre de pages extraites: {len(pages)}")
            return pages
    except Exception as e:
        logger.error(f"Erreur extraction PDF : {str(e)}")
        return f"Erreur lors de l'extraction du texte : {str(e)}"

def extract_pptx_text(filepath):
    try:
        prs = Presentation(filepath)
        pages = []
        for i, slide in enumerate(prs.slides):
            page_text = ''
            sections = []
            # Extraire le titre (si présent)
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    page_text += title_text + '\n'
                    sections.append({
                        'title': title_text,
                        'start': 0,
                        'end': len(title_text)
                    })
            # Extraire le texte des autres formes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip() and shape != slide.shapes.title:
                    page_text += shape.text.strip() + '\n'
            logger.info(f"Diapositive {i+1} extraite, longueur: {len(page_text)} caractères")
            # Détecter d'autres sections dans le texte
            additional_sections = detect_sections(page_text)
            sections.extend(additional_sections)
            pages.append({
                'page_number': i + 1,
                'text': page_text,
                'sections': sections
            })
        if not any(page['text'].strip() for page in pages):
            logger.warning("Aucun texte extrait du PPTX")
            return "Aucun texte lisible n'a été extrait du PPTX."
        logger.info(f"Nombre de diapositives extraites: {len(pages)}")
        return pages
    except Exception as e:
        logger.error(f"Erreur extraction PPTX : {str(e)}")
        return f"Erreur lors de l'extraction du texte : {str(e)}"

def detect_sections(text):
    """
    Détecter les sections avec une regex améliorée et des heuristiques pour les titres.
    """
    sections = []
    # Regex pour chapitres, articles, etc. avec notations arabes et romaines
    section_pattern = r'(?i)\b(chapitre|chapter|article|section|partie|titre|title|module|unité)\s+(\d+|[IVX]+)\b.*?(?=\b(chapitre|chapter|article|section|partie|titre|title|module|unité)\s+(?:\d+|[IVX]+)\b|$|[IVX]+)'
    for match in re.finditer(section_pattern, text):
        sections.append({
            'title': match.group(0).split('\n')[0].strip(),
            'start': match.start(),
            'end': match.end()
        })
    # Heuristique pour titres (ex. "1. INTRODUCTION", lignes en majuscules)
    title_pattern = r'^(?:\d+\.|\b[IVX]+)\s*[A-Z][A-Z\s]{5,50}$|^(?:\d+\.|\b[IVX]+)\s*.+?(?=\n\n|$)'
    for match in re.finditer(title_pattern, text, re.MULTILINE):
        title = match.group(0).strip()
        if title not in [s['title'] for s in sections]:  # Éviter les doublons
            sections.append({
                'title': title,
                'start': match.start(),
                'end': match.end()
            })
    logger.info(f"Sections détectées: {[s['title'] for s in sections]}")
    return sections

def summarize_pages(pages):
    try:
        text = ''.join(page['text'] for page in pages)[:10000]
        summary_prompt = f"""
        Résumez le contenu suivant en 500 mots maximum, en conservant les points clés : {text}
        """
        response = client.chat.completions.create(
            model="meta-llama/Llama-3.3-70B-Instruct-Turbo-Free",
            messages=[{"role": "user", "content": summary_prompt}],
            temperature=0.7,
            max_tokens=750
        )
        summary = response.choices[0].message.content
        logger.info(f"Résumé généré, longueur: {len(summary)} caractères")
        return summary
    except Exception as e:
        logger.error(f"Erreur lors du résumé des pages : {str(e)}")
        return ''.join(page['text'] for page in pages)[:6000]

def get_session_id():
    session_id = request.json.get('session_id') if request.is_json else None
    if not session_id:
        session_id = request.form.get('session_id') or request.headers.get('X-Session-ID', str(uuid.uuid4()))
    logger.info(f"Session ID utilisé: {session_id}")
    return session_id


@app.route('/')
def home():
    logger.info("Accès à la page d'accueil")
    return render_template('homepage.html')

@app.route('/chat.html')
def chat_page():
    logger.info("Accès à la page de chat")
    return render_template('chat.html')

@app.route('/Do-A-Test.html')
def do_a_test():
    logger.info("Accès à la page Do-A-Test")
    return render_template('Do-A-Test.html')

@app.route('/upload_file', methods=['POST'])
def upload_file():
    session_id = get_session_id()
    logger.info(f"Upload fichier pour session {session_id}")
    try:
        if 'file' not in request.files:
            logger.error("Aucun fichier fourni")
            return jsonify({'error': 'Aucun fichier fourni'}), 400
        file = request.files['file']
        is_quiz = request.form.get('is_quiz', 'false').lower() == 'true'
        if file.filename == '':
            logger.error("Aucun fichier sélectionné")
            return jsonify({'error': 'Aucun fichier sélectionné'}), 400
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], f"{session_id}_{filename}")
            file.save(filepath)

            if filename.endswith('.pdf'):
                pages = extract_pdf_text(filepath)
            else:
                pages = extract_pptx_text(filepath)

            if isinstance(pages, str):
                return jsonify({'error': pages}), 500

            with LOCK:
                if session_id not in SESSION_STORAGE:
                    SESSION_STORAGE[session_id] = {'files': [], 'last_file': None}
                SESSION_STORAGE[session_id]['files'].append({
                    'filename': filename,
                    'pages': pages
                })
                SESSION_STORAGE[session_id]['last_file'] = filename
                logger.info(f"Fichiers stockés pour session {session_id}: {[f['filename'] for f in SESSION_STORAGE[session_id]['files']]}, dernier fichier: {filename}")

            os.remove(filepath)
            logger.info(f"Fichier temporaire supprimé: {filepath}")

            # Message en Markdown pour chat.html, texte brut pour Do-A-Test.html
            bot_message = f"Fichier **{filename}** chargé avec succès. Sélectionnez une section pour poser des questions ou générer un quiz." if not is_quiz else f"Fichier {filename} chargé avec succès. Sélectionnez une section pour poser des questions ou générer un quiz."

            logger.info(f"Fichier {filename} chargé, redirection vers {'/Do-A-Test.html' if is_quiz else '/chat.html'}")
            return jsonify({
                'message': bot_message,
                'redirect': '/Do-A-Test.html' if is_quiz else '/chat.html'
            })
        else:
            logger.error("Type de fichier non autorisé")
            return jsonify({'error': 'Type de fichier non autorisé'}), 400
    except Exception as e:
        logger.error(f"Erreur upload fichier : {str(e)}")
        return jsonify({'error': f'Erreur lors du chargement du fichier : {str(e)}'}), 500

@app.route('/chat', methods=['POST'])
def handle_chat():
    session_id = get_session_id()
    logger.info(f"Requête chat pour session {session_id}")
    try:
        data = request.get_json()
        if not data or 'message' not in data:
            logger.error("Le champ 'message' est requis")
            return jsonify({'error': 'Le champ "message" est requis'}), 400
        user_message = data['message']
        user_message_lower = user_message.lower()

        with LOCK:
            session_data = SESSION_STORAGE.get(session_id, {'files': [], 'last_file': None})
            files = session_data['files']
            last_file_filename = session_data.get('last_file')
            logger.info(f"Fichiers stockés pour session {session_id}: {[f['filename'] for f in files]}, dernier fichier: {last_file_filename}")

        target_file = None
        target_filename = None
        for file in files:
            if file['filename'].lower() in user_message_lower:
                target_file = file['pages']
                target_filename = file['filename']
                break
        if not target_file and last_file_filename:
            for file in files:
                if file['filename'] == last_file_filename:
                    target_file = file['pages']
                    target_filename = file['filename']
                    break

        has_file = target_file is not None and isinstance(target_file, list) and any(page['text'].strip() for page in target_file)
        logger.info(f"Fichier cible: {target_filename if target_filename else 'Aucun'}, disponible: {has_file}")

        prompt = ''
        is_quiz = 'quiz' in user_message_lower or 'test' in user_message_lower
        if has_file:
            selected_pages = []
            total_length = 0
            max_length = 6000

            page_range_match = re.search(r'(?:pages?|diapositives?)\s+(\d+)(?:-(\d+))?', user_message_lower)
            if page_range_match:
                start_page = int(page_range_match.group(1))
                end_page = int(page_range_match.group(2)) if page_range_match.group(2) else start_page
                selected_pages = [page for page in target_file if start_page <= page['page_number'] <= end_page]
                logger.info(f"Pages sélectionnées via intervalle: {start_page}-{end_page}")

            section_match = re.search(r'(chapitre|chapter|article|section|partie|titre|title|module|unité)\s+(\d+|[IVX]+)', user_message_lower)
            if section_match and not selected_pages:
                section_title = section_match.group(0)
                for page in target_file:
                    if any(section['title'].lower().startswith(section_title) for section in page['sections']):
                        selected_pages.append(page)
                        next_pages = [p for p in target_file if p['page_number'] > page['page_number']]
                        for next_page in next_pages:
                            if next_page['sections']:
                                break
                            selected_pages.append(next_page)
                        break
                logger.info(f"Pages sélectionnées via section: {section_title}")

            if not selected_pages:
                selected_pages = target_file[:10]
                logger.info("Aucune section ou intervalle spécifié, utilisation des premières pages")

            target_text = ''
            for page in selected_pages:
                if total_length + len(page['text']) <= max_length:
                    target_text += page['text']
                    total_length += len(page['text'])
                else:
                    remaining = max_length - total_length
                    target_text += page['text'][:remaining]
                    total_length = max_length
                    break
            logger.info(f"Texte sélectionné, longueur: {total_length} caractères")

            if total_length > max_length:
                logger.info(f"Texte trop long ({total_length} caractères), génération d'un résumé")
                target_text = summarize_pages(selected_pages)
                total_length = len(target_text)

            # Instructions de formatage : Markdown pour résumés/questions, texte brut pour quiz
            if is_quiz:
                formatting_instructions = """
                Formattez la réponse en texte brut sans utiliser la syntaxe Markdown (ex. évitez *mot*, **mot**, #, etc.).
                Pour les quiz, utilisez le format suivant avec des sauts de ligne entre chaque élément :
                
                Question [numéro] : [Question]
                
                A) [Option]
                
                B) [Option]
                
                C) [Option]
                
                D) [Option]
                
                Ne donne la reponse que si on te le demande.
                
                Laissez une ligne vide entre chaque question.
                Exemple :
                
                Question 1 : Quel est l'objectif principal de la biochimie ?
                
                A) Étudier les propriétés macroscopiques
                
                B) Comprendre les processus chimiques
                
                C) Développer de nouveaux médicaments
                
                D) Étudier les fonctions des organes
                
                
                Question 2 : ...
                """
            else:
                formatting_instructions = """
                Formattez la réponse en Markdown pour un affichage clair.
                - Utilisez **gras** pour les termes importants.
                - Utilisez des listes à puces (`-`) pour les points clés.
                - Séparez les paragraphes par une ligne vide.
                - Évitez les en-têtes (#) sauf si explicitement demandé.
                """

            if is_quiz:
                prompt = f"""
                {formatting_instructions}
                Basé sur le contenu suivant extrait du fichier '{target_filename}' : {target_text},
                générez un quiz de 3 questions à choix multiples (4 options par question) avec les réponses correctes.
                """
            elif any(keyword in user_message_lower for keyword in ['file', 'document', 'paper', 'abstract', 'summary', 'résumé', 'section', 'chapter', 'diapositive']):
                prompt = f"""
                {formatting_instructions}
                Basé sur le contenu suivant extrait du fichier '{target_filename}' : {target_text},
                répondez à la question suivante : {user_message}
                Fournissez une réponse claire et concise basée uniquement sur le contenu du fichier.
                Si la question ne correspond à aucune information du fichier, dites : "Désolé, je n'ai pas trouvé cette information dans le fichier."
                """
            else:
                prompt = f"""
                {formatting_instructions}
                Vous êtes un assistant académique pour les étudiants universitaires.
                Répondez à la question suivante de manière claire et concise : {user_message}
                """
        else:
            formatting_instructions = """
            Formattez la réponse en Markdown pour un affichage clair.
            - Utilisez **gras** pour les termes importants.
            - Utilisez des listes à puces (`-`) pour les points clés.
            - Séparez les paragraphes par une ligne vide.
            - Évitez les en-têtes (#) sauf si explicitement demandé.
            """
            prompt = f"""
            {formatting_instructions}
            Vous êtes un assistant académique pour les étudiants universitaires.
            Répondez à la question suivante de manière claire et concise : {user_message}
            """

        prompt_length = len(prompt)
        max_prompt_length = 7000
        logger.info(f"Taille du prompt: {prompt_length} caractères (~{prompt_length//0.75} tokens)")
        if prompt_length > max_prompt_length:
            logger.warning(f"Prompt trop long ({prompt_length} caractères), réduction du texte")
            excess = prompt_length - max_prompt_length
            target_text = target_text[:len(target_text) - excess]
            prompt = prompt.replace(target_text, target_text[:len(target_text) - excess])
            logger.info(f"Nouvelle taille du prompt: {len(prompt)} caractères (~{len(prompt)//0.75} tokens)")

        logger.info(f"Type de prompt: {'Quiz' if is_quiz else 'Fichier' if has_file and any(keyword in user_message_lower for keyword in ['file', 'document', 'paper', 'abstract', 'summary', 'résumé', 'section', 'chapter', 'diapositive']) else 'Générale'}")
        logger.info(f"Prompt envoyé: {prompt[:200]}...")

        response = client.chat.completions.create(
            model="meta-llama/Llama-3.3-70B-Instruct-Turbo-Free",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1024
        )
        bot_response = response.choices[0].message.content

        logger.info(f"Réponse générée: {bot_response[:200]}...")
        return jsonify({'response': bot_response})
    except Exception as e:
        logger.error(f"Erreur chat : {str(e)}")
        return jsonify({'error': f'Erreur lors du traitement de la requête : {str(e)}'}), 500

@app.route('/get_sections', methods=['POST'])
def get_sections():
    session_id = get_session_id()
    logger.info(f"Demande de sections pour session {session_id}")
    try:
        data = request.get_json()
        if not data or 'filename' not in data:
            logger.error("Le champ 'filename' est requis")
            return jsonify({'error': 'Le champ "filename" est requis'}), 400
        filename = data['filename']

        with LOCK:
            session_data = SESSION_STORAGE.get(session_id, {'files': [], 'last_file': None})
            files = session_data['files']
            target_file = None
            for file in files:
                if file['filename'] == filename:
                    target_file = file['pages']
                    break
            if not target_file:
                logger.error(f"Fichier {filename} non trouvé")
                return jsonify({'error': f'Fichier {filename} non trouvé'}), 404

        sections = []
        for page in target_file:
            for section in page['sections']:
                sections.append({
                    'title': section['title'],
                    'page': page['page_number']
                })
        logger.info(f"Sections trouvées pour {filename}: {len(sections)}")
        return jsonify({'sections': sections})
    except Exception as e:
        logger.error(f"Erreur get_sections : {str(e)}")
        return jsonify({'error': f'Erreur lors de la récupération des sections : {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)